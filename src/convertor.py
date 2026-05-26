"""Beamer PDF → PPTX 转换器。

将 Beamer 生成的 PDF 每页渲染为高 DPI PNG 位图，
插入到等尺寸空白 PPTX 文件中，每页一张图片贴满幻灯片。

依赖: PyMuPDF (fitz), python-pptx, Pillow
"""
import logging
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Callable

logger = logging.getLogger(__name__)

# Beamer 标准宽高比 → PPTX 幻灯片尺寸（英寸）
# python-pptx 用 Inches，这些值对应标准幻灯片尺寸
ASPECT_RATIO_MAP: dict[tuple[int, int], tuple[float, float]] = {
    (4, 3):   (10.00, 7.50),    # 标准 4:3
    (16, 9):  (10.00, 5.625),   # 宽屏 16:9
    (16, 10): (10.00, 6.25),    # 16:10
}

DEFAULT_DPI: int = 250
"""默认渲染 DPI。250 在清晰度和文件体积之间平衡较好。"""


@dataclass
class ConversionResult:
    """转换结果数据容器。"""
    success: bool
    output_path: str | None = None
    error: str | None = None


class BtpConvertor:
    """Beamer to Power Point 转换器。

    将 Beamer 编译输出的 PDF 逐页渲染为高 DPI PNG 图片，
    插入等尺寸空白 PPTX 文件。所有幻灯片为位图格式，
    不支持跳转链接或矢量编辑，但保留了原始排版和字体。

    Attributes:
        source: 源 PDF 路径。
        output_folder: 输出文件夹路径。
        output_format: "PDF to PPTX" 或 "PDF to PPT"（均输出 .pptx）。
        custom_name: 自定义输出文件名（不含扩展名）。
        dpi: 渲染分辨率。
        progress_callback: 进度回调 (current, total) -> None。
    """

    def __init__(
        self,
        source: str,
        output_folder: str,
        output_format: str = "PDF to PPTX",
        custom_name: str = "",
        dpi: int = DEFAULT_DPI,
        progress_callback: Callable[[int, int], None] | None = None,
    ):
        self.source = Path(source)
        self.output_folder = Path(output_folder)
        self.output_format = output_format
        self.custom_name = custom_name
        self.dpi = dpi
        self.progress_callback = progress_callback

    def run(self) -> dict:
        """执行 PDF → PPTX 转换。

        流程:
          1. 用 PyMuPDF 打开 PDF，检查有效性。
          2. 读取首页尺寸，推断 Beamer 宽高比，设置 PPTX 幻灯片尺寸。
          3. 每页渲染为 DPI PNG（内存中），插入空白幻灯片。
          4. 保存 PPTX 到输出路径。

        Returns:
            {"success": True, "output_path": "..."}
            或 {"success": False, "error": "错误描述"}。
        """
        try:
            return self._convert()
        except Exception:
            logger.exception("转换过程发生未预期异常")
            return {"success": False, "error": "转换过程发生未预期异常，请查看日志"}

    # ── 内部实现 ──────────────────────────────────────────

    def _convert(self) -> dict:
        """转换核心逻辑。"""
        import fitz  # PyMuPDF

        # 1. 检查并打开 PDF
        if not self.source.exists():
            return {"success": False, "error": f"文件不存在: {self.source}"}

        doc = fitz.open(str(self.source))
        total_pages = doc.page_count
        if total_pages == 0:
            doc.close()
            return {"success": False, "error": "PDF 文件为空（0 页）"}

        # 2. 推断宽高比
        first_rect = doc[0].rect
        aspect = first_rect.width / first_rect.height
        pptx_width, pptx_height = self._match_aspect_ratio(aspect)

        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(pptx_width)
        prs.slide_height = Inches(pptx_height)

        # 空白布局（python-pptx 索引 6 通常是空白布局）
        blank_layout = prs.slide_layouts[6]

        zoom = self.dpi / 72.0  # PyMuPDF 默认 72 DPI

        # 3. 逐页渲染
        for page_index in range(total_pages):
            if self.progress_callback:
                self.progress_callback(page_index + 1, total_pages)

            page = doc[page_index]
            matrix = fitz.Matrix(zoom, zoom)
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            img_bytes = pixmap.tobytes("png")

            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                BytesIO(img_bytes),
                0, 0,
                Inches(pptx_width),
                Inches(pptx_height),
            )

        doc.close()

        # 4. 确定输出路径并保存
        self.output_folder.mkdir(parents=True, exist_ok=True)

        if self.custom_name:
            out_name = f"{self.custom_name}.pptx"
        else:
            out_name = f"{self.source.stem}_converted.pptx"

        output_path = self.output_folder / out_name
        prs.save(str(output_path))

        logger.info("转换完成: %s -> %s (%d 页, %d DPI)",
                     self.source, output_path, total_pages, self.dpi)

        return {"success": True, "output_path": str(output_path)}

    # ── 静态工具 ──────────────────────────────────────────

    @staticmethod
    def _match_aspect_ratio(aspect: float) -> tuple[float, float]:
        """匹配合适的 Beamer 宽高比，返回 (宽英寸, 高英寸)。

        从 ASPECT_RATIO_MAP 中找到与输入宽高比最接近的标准比例。
        """
        best_key = min(
            ASPECT_RATIO_MAP.keys(),
            key=lambda k: abs(aspect - k[0] / k[1]),
        )
        return ASPECT_RATIO_MAP[best_key]
